import traceback
from pathlib import Path

# 导入增强模块
try:
    from src.core.data_masking import mask_for_bug_analysis, mask_for_table_check, DataMasker
    from src.core.bug_analyzer import BugAnalyzer
    from src.core.table_checker import TableChecker
    ENHANCED_MODULES = True
except ImportError:
    ENHANCED_MODULES = False

# 智能模型选择
def get_model_for_task(task_type, text_len):
    """根据任务类型和复杂度选择模型"""
    if task_type == 'chat':
        return 'kimi-k2-turbo-preview'  # 对话用快速模型
    elif task_type == 'check':
        return 'kimi-k2-turbo-preview'  # 表检查用快速模型
    elif task_type == 'analyze':
        if text_len > 3000:
            return 'kimi-k2-pro'  # 长文档用强模型
        return 'kimi-k2-turbo-preview'
    elif task_type == 'generate':
        return 'kimi-k2-pro'  # 用例生成需要高质量
    return 'kimi-k2-turbo-preview'

import re
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import os
import sys
import secrets
import time
import tempfile
import shutil
from typing import Optional
import asyncio

# 动态添加项目根目录到路径
PROJECT_ROOT = Path(__file__).parent.parent.resolve()
sys.path.insert(0, str(PROJECT_ROOT))

from src.core.config import Config

app = FastAPI(title="TestOwl API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

sessions = {}

# 存储上传文件内容的字典（内存存储，重启后丢失）
file_contents = {}

class SaveKeyRequest(BaseModel):
    api_key: str

class AnalyzeRequest(BaseModel):
    text: str
    session_token: Optional[str] = None

class UrlRequest(BaseModel):
    url: str

def get_api_key(session_token: Optional[str]) -> Optional[str]:
    if session_token and session_token in sessions:
        if time.time() - sessions[session_token]['created'] < 7 * 24 * 3600:
            return sessions[session_token]['api_key']
        else:
            del sessions[session_token]
    return None

def get_config_with_key(api_key: str):
    """创建带有指定API key的config"""
    config = Config('/root/testowl/config/config.yaml')
    config.llm.api_key = api_key
    return config

def parse_file(file_path: str, suffix: str) -> str:
    """解析各种格式的文件内容为文本"""
    content = ""
    try:
        if suffix == '.txt' or suffix == '.md':
            # 文本文件和 Markdown 文件
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        elif suffix in ['.docx', '.doc']:
            # Word 文档
            try:
                from docx import Document
                doc = Document(file_path)
                content = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            except ImportError:
                return "[错误：未安装 python-docx，请运行：pip install python-docx]"
            except Exception as e:
                return f"[Word文档解析失败: {str(e)}]"
        elif suffix == '.pdf':
            # PDF 文件
            try:
                # 尝试使用 pypdf (新版)
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                content = '\n'.join([page.extract_text() or '' for page in reader.pages])
            except ImportError:
                # 回退到 PyPDF2 (旧版)
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(file_path)
                    content = '\n'.join([page.extract_text() or '' for page in reader.pages])
                except ImportError:
                    return "[错误：未安装 PDF 解析库，请运行：pip install pypdf 或 pip install PyPDF2]"
            except Exception as e:
                return f"[PDF解析失败: {str(e)}]"
        elif suffix in ['.xlsx', '.xls']:
            # Excel 文件
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, data_only=True)
                texts = []
                for sheet in wb.worksheets:
                    texts.append(f"=== {sheet.title} ===")
                    for row in sheet.iter_rows(values_only=True):
                        if any(cell for cell in row):
                            texts.append(' | '.join(str(cell or '') for cell in row))
                content = '\n'.join(texts)
            except ImportError:
                return "[错误：未安装 openpyxl，请运行：pip install openpyxl]"
            except Exception as e:
                return f"[Excel解析失败: {str(e)}]"
        elif suffix == '.pptx':
            # PPT 文件
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                texts = []
                for i, slide in enumerate(prs.slides):
                    texts.append(f"=== 第{i+1}页 ===")
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            texts.append(shape.text)
                content = '\n'.join(texts)
            except ImportError:
                return "[错误：未安装 python-pptx，请运行：pip install python-pptx]"
            except Exception as e:
                return f"[PPT解析失败: {str(e)}]"
        elif suffix in ['.png', '.jpg', '.jpeg', '.bmp', '.gif', '.webp']:
            # 图片文件 - OCR
            try:
                import pytesseract
                from PIL import Image
                img = Image.open(file_path)
                content = pytesseract.image_to_string(img, lang='chi_sim+eng')
                if not content.strip():
                    content = "[图片中未识别到文字内容]"
            except ImportError:
                return "[错误：OCR功能未安装，请运行：pip install pytesseract Pillow，并安装Tesseract-OCR引擎]"
            except Exception as e:
                return f"[图片OCR失败: {str(e)}]"
        elif suffix in ['.csv']:
            # CSV 文件
            try:
                import csv
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = [' | '.join(row) for row in reader]
                    content = '\n'.join(rows)
            except Exception as e:
                return f"[CSV解析失败: {str(e)}]"
        elif suffix in ['.json']:
            # JSON 文件
            try:
                import json
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    data = json.load(f)
                    content = json.dumps(data, ensure_ascii=False, indent=2)
            except Exception as e:
                return f"[JSON解析失败: {str(e)}]"
        else:
            # 不支持的格式
            return f"[不支持的文件格式: {suffix}，支持的格式：txt, md, docx, pdf, xlsx, pptx, png, jpg, csv, json]"
    except Exception as e:
        return f"[文件解析错误: {str(e)}]"
    
    return content.strip()

def fetch_url(url: str) -> str:
    try:
        from bs4 import BeautifulSoup
        import requests
        headers = {'User-Agent': 'Mozilla/5.0'}
        resp = requests.get(url, headers=headers, timeout=15)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, 'html.parser')
        for tag in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
            tag.decompose()
        title = soup.title.string if soup.title else ''
        text = soup.get_text(separator='\n', strip=True)
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        return f"标题: {title}\n\n" + '\n'.join(lines)
    except requests.RequestException as e:
        raise RuntimeError(f"网络请求失败: {str(e)}")
    except Exception as e:
        raise RuntimeError(f"获取网页内容失败: {str(e)}")

@app.post("/auth/save_key")
async def save_key(req: SaveKeyRequest):
    token = secrets.token_urlsafe(32)
    sessions[token] = {'api_key': req.api_key, 'created': time.time()}
    return {"success": True, "session_token": token}

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """上传文件并解析内容"""
    try:
        # 确保上传目录存在
        upload_dir = Path("uploads")
        upload_dir.mkdir(exist_ok=True)
        
        # 生成唯一文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        suffix = os.path.splitext(file.filename)[1].lower()
        safe_filename = f"{timestamp}_{file.filename}"
        file_path = upload_dir / safe_filename
        
        # 保存文件
        with open(file_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        
        # 解析文件内容
        content = parse_file(str(file_path), suffix)
        
        # 生成文件ID
        file_id = f"file_{timestamp}"
        
        # 保存文件内容到内存（供后续分析使用）
        file_contents[file_id] = {
            "filename": file.filename,
            "content": content,
            "path": str(file_path),
            "upload_time": datetime.now().isoformat()
        }
        
        # 清理临时文件（保留原始文件）
        # os.unlink(tmp_path)
        
        if not content:
            return {"success": False, "error": "文件内容为空或无法解析"}
        
        return {
            "success": True, 
            "file_id": file_id,
            "filename": file.filename,
            "text": content[:500] + "..." if len(content) > 500 else content  # 返回前500字符预览
        }
    except Exception as e:
        import traceback
        return {"success": False, "error": str(e), "traceback": traceback.format_exc()}

@app.post("/fetch_url")
async def fetch_url_endpoint(req: UrlRequest):
    try:
        content = fetch_url(req.url)
        return {"success": True, "text": content}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/file/{file_id}")
async def get_file(file_id: str):
    """获取上传文件的内容"""
    if file_id not in file_contents:
        return {"success": False, "error": "文件不存在或已过期"}
    
    file_info = file_contents[file_id]
    return {
        "success": True,
        "file_id": file_id,
        "filename": file_info["filename"],
        "content": file_info["content"]
    }

@app.post("/analyze")
async def analyze(req: AnalyzeRequest):
    key = get_api_key(req.session_token)
    if not key:
        return {"success": False, "error": "未配置API Key"}
    try:
        config = get_config_with_key(key)
        from src.adapters.llm.client import LLMClient
        
        text = req.text[:6000] if len(req.text) > 6000 else req.text
        model = get_model_for_task('analyze', len(text))
        config.llm.model = model
        client = LLMClient(config)
        
        prompt = "分析需求，提取测试点、风险点、待确认问题。返回JSON: " + text + " 必须包含: {document_summary,test_points:[{id,title,description,priority,category}],risk_points:[{description,impact}],questions:[{question}]}"
        result = await asyncio.wait_for(client.complete(prompt), timeout=90)
        
        import json
        clean = result.strip()
        if "```" in clean:
            parts = clean.split("```")
            clean = parts[1] if len(parts) > 1 else parts[0]
            clean = clean.replace("json", "", 1).strip()
        
        return {"success": True, "data": json.loads(clean)}
    except asyncio.TimeoutError:
        return {"success": False, "error": "分析超时"}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.post("/generate_cases")
async def generate_cases(req: AnalyzeRequest):
    key = get_api_key(req.session_token)
    if not key:
        return {"success": False, "error": "未配置API Key"}
    try:
        config = get_config_with_key(key)
        from src.adapters.llm.client import LLMClient
        
        text = req.text[:5000] if len(req.text) > 5000 else req.text
        model = get_model_for_task('generate', len(text))
        config.llm.model = model
        client = LLMClient(config)
        
        prompt = f"根据以下需求生成测试用例JSON: {text} 格式: test_cases数组"
        result = await asyncio.wait_for(client.complete(prompt), timeout=90)
        
        import json
        clean_result = result.strip()
        if "```" in clean_result:
            clean_result = clean_result.split("```")[1]
            if clean_result.startswith("json"):
                clean_result = clean_result[4:]
        
        return {"success": True, "data": json.loads(clean_result)}
    except asyncio.TimeoutError:
        return {"success": False, "error": "生成超时"}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.get("/health")
async def health():
    """真实冒烟测试"""
    import time, os, traceback
    from io import BytesIO
    from datetime import datetime
    
    checks = []
    
    # 1. API服务检查
    checks.append({"name": "API服务", "status": "pass", "message": "运行正常"})
    
    # 2. 会话管理检查
    try:
        active_count = len(sessions)
        checks.append({"name": "会话管理", "status": "pass", "message": f"正常，{active_count}个活跃会话"})
    except Exception as e:
        checks.append({"name": "会话管理", "status": "fail", "message": str(e)})
    
    # 3. 文件存储检查
    try:
        upload_dir = "uploads"
        if not os.path.exists(upload_dir):
            os.makedirs(upload_dir)
        file_count = len(os.listdir(upload_dir)) if os.path.exists(upload_dir) else 0
        checks.append({"name": "文件存储", "status": "pass", "message": f"目录正常，{file_count}个文件"})
    except Exception as e:
        checks.append({"name": "文件存储", "status": "fail", "message": str(e)})
    
    # 4. 导出功能检查 - 模拟前端实际请求参数
    export_formats = [
        ("md", "Markdown"),
        ("pdf", "PDF"),
        ("xlsx", "Excel"),  # 前端实际传的是xlsx
        ("docx", "Word")    # 前端实际传的是docx
    ]
    
    for fmt, name in export_formats:
        try:
            # 真实生成文件到内存
            if fmt == "md":
                file_bytes = "测试内容".encode('utf-8')
            elif fmt == "pdf":
                from reportlab.lib.pagesizes import A4
                from reportlab.pdfgen import canvas
                from reportlab.pdfbase import pdfmetrics
                from reportlab.pdfbase.ttfonts import TTFont
                try:
                    pdfmetrics.registerFont(TTFont('SimSun', '/usr/share/fonts/truetype/wqy/wqy-microhei.ttc'))
                    font = 'SimSun'
                except Exception:
                    font = 'Helvetica'
                buffer = BytesIO()
                c = canvas.Canvas(buffer, pagesize=A4)
                c.setFont(font, 10)
                c.drawString(50, 800, "测试内容")
                c.save()
                file_bytes = buffer.getvalue()
            elif fmt == "xlsx":
                from openpyxl import Workbook
                wb = Workbook()
                ws = wb.active
                ws.cell(row=1, column=1, value="测试内容")
                buffer = BytesIO()
                wb.save(buffer)
                file_bytes = buffer.getvalue()
            elif fmt == "docx":
                from docx import Document
                doc = Document()
                doc.add_paragraph("测试内容")
                buffer = BytesIO()
                doc.save(buffer)
                file_bytes = buffer.getvalue()
            
            checks.append({"name": f"{name}导出", "status": "pass", "message": "生成成功"})
        except Exception as e:
            tb = traceback.format_exc()
            checks.append({
                "name": f"{name}导出", 
                "status": "fail", 
                "message": str(e)[:50],
                "traceback": tb
            })
    
    # 5. 配置目录检查
    try:
        config_exists = os.path.exists("config")
        checks.append({"name": "配置目录", "status": "pass" if config_exists else "warn", "message": "存在" if config_exists else "不存在"})
    except Exception as e:
        checks.append({"name": "配置目录", "status": "fail", "message": str(e)})
    
    # 6. 导出功能可用性检查
    try:
        # 简单检查前端是否有导出相关函数
        with open('/root/testowl/web/index.html', 'r') as f:
            html = f.read()
        has_export = 'exportBubble' in html or 'exportChat' in html
        if has_export:
            checks.append({"name": "导出功能", "status": "pass", "message": "导出功能已启用"})
        else:
            checks.append({"name": "导出功能", "status": "warn", "message": "导出功能未找到"})
    except Exception as e:
        checks.append({"name": "导出功能", "status": "warn", "message": f"检查失败: {str(e)[:30]}"})
    
    return {
        "status": "ok" if all(c["status"] == "pass" for c in checks) else "warn",
        "timestamp": datetime.now().isoformat(),
        "checks": checks
    }


@app.post("/export")
async def export_report(req: AnalyzeRequest):
    key = get_api_key(req.session_token)
    if not key:
        return {"success": False, "error": "未配置API Key"}
    try:
        config = get_config_with_key(key)
        from src.adapters.llm.client import LLMClient
        client = LLMClient(config)
        
        prompt = "生成测试分析报告(Markdown格式): " + req.text[:3000]
        result = await asyncio.wait_for(client.complete(prompt), timeout=60)
        
        from datetime import datetime
        return {
            "success": True,
            "filename": f"test_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md",
            "content": result
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


from web.chat_handler import handle_chat, detect_intent

@app.post("/chat")
async def chat(req: dict):
    """处理聊天请求，支持增强的报错分析和表检查"""
    user_msg = req.get("message", "")
    token = req.get("session_token", "")
    file_id = req.get("file_id", "")
    
    key = get_api_key(token)
    if not key:
        return {"success": False, "error": "请先配置API Key"}
    
    try:
        # 检测意图
        config = get_config_with_key(key)
        from src.adapters.llm.client import LLMClient
        client = LLMClient(config)
        
        intent_data = await detect_intent(client, user_msg, [])
        intent = intent_data.get("intent", "chat")
        
        # 报错分析 - 使用增强模块
        if intent == "analyze_bug" and ENHANCED_MODULES:
            # 检查是否有文件内容（错误日志）
            error_text = ""
            if file_id and file_id in file_contents:
                error_text = file_contents[file_id].get("content", "")
                print(f"[DEBUG] Got file content, length: {len(error_text)}, preview: {error_text[:100]}")
            
            # 检查是否是OCR失败的提示
            if error_text and ("[图片中未识别到文字" in error_text or "[OCR失败" in error_text or "[错误：" in error_text):
                return {"success": True, "response": f"无法识别图片内容：{error_text}<br><br>建议直接粘贴错误日志文本，或确保截图清晰包含错误信息。"}
            
            # 如果没有文件内容或内容太短，尝试从消息中提取
            if not error_text or len(error_text) < 20:
                # 清理消息中的指令性文字，保留可能的错误内容
                cleaned_msg = re.sub(r'^(分析|查看|检查|帮忙|请|帮我)\s*', '', user_msg)
                if len(cleaned_msg) > len(error_text):
                    error_text = cleaned_msg
            
            if error_text and len(error_text) > 10:
                # 脱敏处理
                masked_text, mask_records = mask_for_bug_analysis(error_text)
                
                # 使用增强分析器
                analyzer = BugAnalyzer()
                report = analyzer.analyze(masked_text)
                html_report = analyzer.generate_html_report(report)
                
                # 添加脱敏说明
                if mask_records:
                    html_report += f"<p style='color:#999;font-size:12px;margin-top:10px;'>[安全] 已自动脱敏 {len(mask_records)} 处敏感信息</p>"
                
                return {"success": True, "response": html_report}
            else:
                # 没有足够的内容进行分析
                return {"success": True, "response": "请提供报错日志或截图，我可以帮你分析错误原因。你可以直接粘贴错误信息或上传日志文件。"}
        
        # 表检查 - 使用增强模块
        elif intent == "check_table" and ENHANCED_MODULES:
            # 检查是否有文件内容（表格数据）
            if file_id and file_id in file_contents:
                file_info = file_contents[file_id]
                content = file_info.get("content", "")
                filename = file_info.get("filename", "")
                
                # 尝试解析表格数据
                table_data = parse_table_from_text(content)
                
                if table_data and table_data.get("rows"):
                    # 脱敏处理
                    for row in table_data.get("rows", []):
                        for key, val in row.items():
                            if isinstance(val, str):
                                masked_val, _ = mask_for_table_check(val)
                                row[key] = masked_val
                    
                    # 使用增强检查器
                    checker = TableChecker()
                    report = checker.check(table_data, filename)
                    html_report = checker.generate_html_report(report)
                    
                    return {"success": True, "response": html_report}
                else:
                    return {"success": True, "response": "无法解析表格数据，请确保上传的是有效的Excel/CSV文件。"}
        
        # 默认使用原有的handle_chat
        return await handle_chat(req, get_api_key, get_config_with_key, file_contents)
    
    except Exception as e:
        # 出错时回退到原有处理
        return await handle_chat(req, get_api_key, get_config_with_key, file_contents)


def parse_table_from_text(text: str) -> dict:
    """从文本中解析表格数据"""
    if not text:
        return None
    
    lines = text.strip().split('\n')
    if len(lines) < 2:
        return None
    
    # 尝试检测分隔符
    # 1. 尝试制表符
    if '\t' in lines[0]:
        headers = lines[0].split('\t')
        rows = []
        for line in lines[1:]:
            if line.strip():
                values = line.split('\t')
                row = {headers[i]: values[i] if i < len(values) else '' for i in range(len(headers))}
                rows.append(row)
        return {'headers': headers, 'rows': rows}
    
    # 2. 尝试逗号（CSV）
    elif ',' in lines[0]:
        headers = lines[0].split(',')
        rows = []
        for line in lines[1:]:
            if line.strip():
                values = line.split(',')
                row = {headers[i]: values[i] if i < len(values) else '' for i in range(len(headers))}
                rows.append(row)
        return {'headers': headers, 'rows': rows}
    
    # 3. 尝试空格（对齐）
    else:
        # 简单处理：第一行是标题，后面是数据
        headers = lines[0].split()
        rows = []
        for line in lines[1:]:
            if line.strip():
                values = line.split()
                row = {}
                for i, h in enumerate(headers):
                    row[h] = values[i] if i < len(values) else ''
                rows.append(row)
        return {'headers': headers, 'rows': rows}

@app.post("/export_chat")
async def export_chat(req: dict):
    """导出对话记录为报告"""
    token = req.get("session_token")
    messages = req.get("messages", [])
    fmt = req.get("format", "md")
    
    key = get_api_key(token)
    if not key:
        return {"success": False, "error": "未配置API Key"}
    
    try:
        config = get_config_with_key(key)
        from src.adapters.llm.client import LLMClient
        client = LLMClient(config)
        
        # 构建对话内容
        chat_text = "\n".join([f"{m['role']}: {m['content']}" for m in messages])
        
        if fmt == "html":
            prompt = f"将以下对话整理为HTML格式的测试报告，包含样式：\n{chat_text}"
        elif fmt == "txt":
            prompt = f"将以下对话整理为纯文本格式的测试报告：\n{chat_text}"
        else:
            prompt = f"将以下对话整理为Markdown格式的测试报告，包含标题、测试要点、结论：\n{chat_text}"
        
        result = await asyncio.wait_for(client.complete(prompt), timeout=60)
        
        from datetime import datetime
        ext = {"md": "md", "txt": "txt", "html": "html"}.get(fmt, "md")
        
        return {
            "success": True,
            "filename": f"chat_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}",
            "content": result
        }
    except Exception as e:
        return {"success": False, "error": str(e)}

import base64
from datetime import datetime
from io import BytesIO

def html_to_text(html: str) -> str:
    """将 HTML 转换为纯文本，保留结构信息"""
    import re
    
    # 移除 script 和 style 标签及其内容
    html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
    html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
    
    # 将标题标签转换为 Markdown 格式
    html = re.sub(r'<h1[^>]*>(.*?)</h1>', r'# \1\n', html, flags=re.IGNORECASE)
    html = re.sub(r'<h2[^>]*>(.*?)</h2>', r'## \1\n', html, flags=re.IGNORECASE)
    html = re.sub(r'<h3[^>]*>(.*?)</h3>', r'### \1\n', html, flags=re.IGNORECASE)
    html = re.sub(r'<h4[^>]*>(.*?)</h4>', r'#### \1\n', html, flags=re.IGNORECASE)
    
    # 将列表项转换为 Markdown 格式
    html = re.sub(r'<li[^>]*>(.*?)</li>', r'- \1\n', html, flags=re.IGNORECASE)
    
    # 将段落转换为文本
    html = re.sub(r'<p[^>]*>(.*?)</p>', r'\1\n\n', html, flags=re.IGNORECASE | re.DOTALL)
    
    # 将 <br> 转换为换行
    html = re.sub(r'<br\s*/?>', '\n', html, flags=re.IGNORECASE)
    
    # 移除所有其他 HTML 标签
    html = re.sub(r'<[^>]+>', '', html)
    
    # 解码 HTML 实体
    import html as html_module
    html = html_module.unescape(html)
    
    # 清理多余空白
    lines = html.split('\n')
    cleaned_lines = []
    for line in lines:
        line = line.strip()
        if line:
            cleaned_lines.append(line)
    
    return '\n'.join(cleaned_lines)


def parse_analysis_content(content: str) -> dict:
    """解析需求分析内容，提取结构化数据
    
    支持两种输入格式：
    1. HTML 格式（来自 chat_handler）
    2. 纯文本/Markdown 格式
    """
    import re
    
    result = {
        "title": "需求分析报告",
        "summary": "",
        "test_points": [],
        "risks": []
    }
    
    # 如果是 HTML，先提取文本内容
    if content.strip().startswith('<') and '</' in content:
        content = html_to_text(content)
    
    lines = content.split('\n')
    current_section = None
    summary_lines = []
    
    for line in lines:
        line = line.strip()
        if not line:
            # 空行处理：如果在概述部分，保留一个换行
            if current_section == "summary" and summary_lines:
                summary_lines.append("")
            continue
            
        # 检测标题（支持各种格式：# 需求分析、### 需求分析、需求分析等）
        if re.match(r'^#*\s*需求分析', line) or re.match(r'^#*\s*测试分析', line):
            result["title"] = re.sub(r'^#+\s*', '', line)  # 移除 Markdown 标记
            current_section = "summary"
        elif re.match(r'^#*\s*测试点', line):
            # 保存已收集的概述
            if summary_lines:
                result["summary"] = '\n'.join(summary_lines).strip()
            current_section = "test_points"
        elif re.match(r'^#*\s*风险', line):
            # 保存已收集的概述
            if summary_lines and not result["summary"]:
                result["summary"] = '\n'.join(summary_lines).strip()
            current_section = "risks"
        elif current_section == "summary":
            # 收集概述内容（多行）
            summary_lines.append(line)
        elif current_section == "test_points":
            # 解析测试点 [P0] 标题 - 描述
            if line.startswith('•') or line.startswith('-') or line.startswith('*') or re.match(r'^\d+\.', line):
                # 提取优先级和内容
                priority_match = re.search(r'\[([Pp]\d+)\]', line)
                priority = priority_match.group(1).upper() if priority_match else "P2"
                
                # 移除 bullet、序号和优先级标记
                clean_line = re.sub(r'^[•\-\*]\s*', '', line)
                clean_line = re.sub(r'^\d+\.\s*', '', clean_line)
                clean_line = re.sub(r'\[[Pp]\d+\]\s*', '', clean_line)
                
                # 分割标题和描述
                if ' - ' in clean_line:
                    title, desc = clean_line.split(' - ', 1)
                elif '：' in clean_line:
                    title, desc = clean_line.split('：', 1)
                elif ':' in clean_line:
                    title, desc = clean_line.split(':', 1)
                else:
                    title, desc = clean_line, ""
                
                result["test_points"].append({
                    "priority": priority,
                    "title": title.strip(),
                    "description": desc.strip()
                })
        elif current_section == "risks":
            if line.startswith('•') or line.startswith('-') or line.startswith('*') or re.match(r'^\d+\.', line):
                risk = re.sub(r'^[•\-\*]\s*', '', line)
                risk = re.sub(r'^\d+\.\s*', '', risk)
                result["risks"].append(risk)
    
    # 最后保存概述（如果还没有保存）
    if summary_lines and not result["summary"]:
        result["summary"] = '\n'.join(summary_lines).strip()
    
    return result


def parse_testcase_content(content: str) -> dict:
    """解析测试用例内容，提取结构化数据（支持HTML格式）"""
    result = {
        "title": "测试用例报告",
        "cases": []
    }
    
    # 首先尝试从HTML中提取JSON原始数据
    # 如果内容包含JSON格式，直接解析
    try:
        # 查找可能的JSON数组
        json_match = re.search(r'"cases"\s*:\s*(\[.*?\])', content, re.DOTALL)
        if json_match:
            json_str = '{"cases":' + json_match.group(1) + '}'
            # 清理HTML标签
            json_str = re.sub(r'<[^>]+>', '', json_str)
            # 修复HTML实体
            json_str = json_str.replace('&lt;', '<').replace('&gt;', '>').replace('&amp;', '&')
            data = json.loads(json_str)
            cases = []
            for c in data.get('cases', []):
                cases.append({
                    "id": c.get('id', 'TC001'),
                    "title": c.get('title', ''),
                    "priority": c.get('pri', 'P2'),
                    "steps": c.get('steps', []),
                    "expected": c.get('expected', '')
                })
            if cases:
                result["cases"] = cases
                return result
    except Exception:
        pass
    
    # 回退到HTML解析
    # 清理HTML标签，提取纯文本
    clean_content = re.sub(r'<[^>]+>', '\n', content)
    # 修复HTML实体
    clean_content = clean_content.replace('&lt;', '<').replace('&gt;', '>').replace('&amp;', '&')
    
    lines = clean_content.split('\n')
    current_case = None
    in_steps = False
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # 检测标题
        if line.startswith('测试用例'):
            result["title"] = line
        # 检测用例标题行 (TC001 xxx [P0]) - 支持Markdown格式 **TC001**
        elif re.match(r'\*?\*?TC\d+', line):
            # 保存上一个用例
            if current_case:
                result["cases"].append(current_case)
            
            # 清理Markdown标记 **
            clean_line = re.sub(r'\*+', '', line).strip()
            
            # 解析新用例 - 支持多种格式
            # 格式1: TC001 使用正确的用户名和密码成功登录 [P0]
            # 格式2: TC001 使用正确的用户名和密码成功登录[P0]
            match = re.match(r'(TC\d+)\s+(.+?)\s*\[([Pp]\d+)\]', clean_line)
            if match:
                case_id = match.group(1)
                title = match.group(2).strip()
                priority = match.group(3).upper()
            else:
                # 尝试没有空格的格式
                match2 = re.match(r'(TC\d+)\s+(.+?)\[([Pp]\d+)\]', clean_line)
                if match2:
                    case_id = match2.group(1)
                    title = match2.group(2).strip()
                    priority = match2.group(3).upper()
                else:
                    # 最简单的格式：只有TC编号
                    case_id = re.match(r'(TC\d+)', clean_line).group(1)
                    title = clean_line.replace(case_id, '').strip()
                    priority = "P2"
            
            current_case = {
                "id": case_id,
                "title": title,
                "priority": priority,
                "steps": [],
                "expected": ""
            }
            in_steps = True
        # 检测预期结果行（以"预期:"开头）
        elif current_case and (line.startswith('预期:') or line.startswith('预期：')):
            current_case["expected"] = line.replace('预期:', '').replace('预期：', '').strip()
            in_steps = False
        # 检测包含"预期"关键词的行
        elif current_case and ('预期' in line and not line.startswith('TC')):
            # 如果行以"预期"开头或在中间
            if line.startswith('预期'):
                current_case["expected"] = re.sub(r'^预期[：:]\s*', '', line)
            else:
                # 可能是 "预期:xxx" 格式
                expected_match = re.search(r'预期[：:]\s*(.+)', line)
                if expected_match:
                    current_case["expected"] = expected_match.group(1).strip()
            in_steps = False
        # 检测步骤（在用例标题之后，预期结果之前的行）
        elif current_case and in_steps and not line.startswith('TC'):
            # 这可能是步骤内容
            # 移除可能的序号前缀
            step = re.sub(r'^\d+[\.、]\s*', '', line)
            if step and not step.startswith('预期'):
                current_case["steps"].append(step)
    
    # 保存最后一个用例
    if current_case:
        result["cases"].append(current_case)
    
    return result


def export_analysis_report(parsed: dict, fmt: str, timestamp: str) -> dict:
    """导出需求分析报告"""
    try:
        if fmt == "md":
            md_lines = [f"# {parsed['title']}", ""]
            
            if parsed['summary']:
                md_lines.extend(["## 概述", parsed['summary'], ""])
            
            if parsed['test_points']:
                md_lines.extend(["## 测试点", ""])
                for tp in parsed['test_points']:
                    md_lines.append(f"### [{tp['priority']}] {tp['title']}")
                    if tp['description']:
                        md_lines.append(f"{tp['description']}")
                    md_lines.append("")
            
            if parsed['risks']:
                md_lines.extend(["## 风险", ""])
                for risk in parsed['risks']:
                    md_lines.append(f"- {risk}")
                md_lines.append("")
            
            file_bytes = '\n'.join(md_lines).encode('utf-8')
            filename = f"analysis_report_{timestamp}.md"
        
        elif fmt == "pdf":
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib import colors
            
            font_name = register_pdf_font()
            
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4,
                                  rightMargin=72, leftMargin=72,
                                  topMargin=72, bottomMargin=18)
            
            styles = getSampleStyleSheet()
            styles.add(ParagraphStyle(name='ChineseTitle', fontName=font_name, fontSize=18, spaceAfter=20))
            styles.add(ParagraphStyle(name='ChineseHeading', fontName=font_name, fontSize=14, spaceAfter=12, spaceBefore=12))
            styles.add(ParagraphStyle(name='ChineseBody', fontName=font_name, fontSize=10, spaceAfter=6))
            
            story = []
            
            story.append(Paragraph(parsed['title'], styles['ChineseTitle']))
            story.append(Spacer(1, 0.2*inch))
            
            if parsed['summary']:
                story.append(Paragraph("概述", styles['ChineseHeading']))
                story.append(Paragraph(parsed['summary'], styles['ChineseBody']))
                story.append(Spacer(1, 0.1*inch))
            
            if parsed['test_points']:
                story.append(Paragraph("测试点", styles['ChineseHeading']))
                story.append(Spacer(1, 0.1*inch))
                
                data = [['优先级', '测试项', '描述']]
                for tp in parsed['test_points']:
                    data.append([tp['priority'], tp['title'], tp['description'][:50] + '...' if len(tp['description']) > 50 else tp['description']])
                
                table = Table(data, colWidths=[0.8*inch, 2*inch, 3*inch])
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#c4a77d')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, -1), font_name),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                    ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f5f0e8')]),
                ]))
                story.append(table)
                story.append(Spacer(1, 0.2*inch))
            
            if parsed['risks']:
                story.append(Paragraph("风险", styles['ChineseHeading']))
                for risk in parsed['risks']:
                    story.append(Paragraph(f"• {risk}", styles['ChineseBody']))
            
            doc.build(story)
            file_bytes = buffer.getvalue()
            filename = f"analysis_report_{timestamp}.pdf"
        
        elif fmt == "xlsx":
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            
            wb = Workbook()
            
            if parsed['test_points']:
                ws = wb.active
                ws.title = "测试点"
                
                headers = ['优先级', '测试项', '描述']
                ws.append(headers)
                
                header_fill = PatternFill(start_color='C4A77D', end_color='C4A77D', fill_type='solid')
                header_font = Font(bold=True, color='FFFFFF')
                for cell in ws[1]:
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = Alignment(horizontal='center', vertical='center')
                
                for tp in parsed['test_points']:
                    ws.append([tp['priority'], tp['title'], tp['description']])
                
                ws.column_dimensions['A'].width = 10
                ws.column_dimensions['B'].width = 30
                ws.column_dimensions['C'].width = 50
            
            if parsed['risks']:
                ws_risk = wb.create_sheet("风险")
                ws_risk.append(['序号', '风险描述'])
                
                header_fill = PatternFill(start_color='C4A77D', end_color='C4A77D', fill_type='solid')
                header_font = Font(bold=True, color='FFFFFF')
                for cell in ws_risk[1]:
                    cell.fill = header_fill
                    cell.font = header_font
                
                for i, risk in enumerate(parsed['risks'], 1):
                    ws_risk.append([i, risk])
                
                ws_risk.column_dimensions['A'].width = 8
                ws_risk.column_dimensions['B'].width = 80
            
            buffer = BytesIO()
            wb.save(buffer)
            file_bytes = buffer.getvalue()
            filename = f"analysis_report_{timestamp}.xlsx"
        
        elif fmt == "docx":
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
            
            doc = Document()
            
            # 设置文档默认字体为通用字体
            def set_run_font(run, font_name='Arial', font_size=11, bold=False):
                run.font.name = font_name
                run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                run.font.size = Pt(font_size)
                run.bold = bold
            
            # 添加标题
            title = doc.add_heading(level=0)
            title_run = title.add_run(parsed['title'])
            set_run_font(title_run, font_size=18, bold=True)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            if parsed['summary']:
                h = doc.add_heading(level=1)
                set_run_font(h.add_run('概述'), font_size=14, bold=True)
                p = doc.add_paragraph()
                set_run_font(p.add_run(parsed['summary']), font_size=11)
            
            if parsed['test_points']:
                h = doc.add_heading(level=1)
                set_run_font(h.add_run('测试点'), font_size=14, bold=True)
                for tp in parsed['test_points']:
                    p = doc.add_paragraph()
                    set_run_font(p.add_run(f"[{tp['priority']}] "), font_size=11, bold=True)
                    set_run_font(p.add_run(tp['title']), font_size=11, bold=True)
                    if tp['description']:
                        p2 = doc.add_paragraph(style='List Bullet')
                        set_run_font(p2.add_run(tp['description']), font_size=10)
            
            if parsed['risks']:
                h = doc.add_heading(level=1)
                set_run_font(h.add_run('风险'), font_size=14, bold=True)
                for risk in parsed['risks']:
                    p = doc.add_paragraph(style='List Bullet')
                    set_run_font(p.add_run(risk), font_size=10)
            
            buffer = BytesIO()
            doc.save(buffer)
            file_bytes = buffer.getvalue()
            filename = f"analysis_report_{timestamp}.docx"
        
        else:
            return {"success": False, "error": f"不支持的格式: {fmt}"}
        
        return {
            "success": True,
            "file": base64.b64encode(file_bytes).decode('utf-8'),
            "filename": filename
        }
    
    except Exception as e:
        import traceback
        return {"success": False, "error": f"导出失败: {str(e)}", "traceback": traceback.format_exc()}


def export_testcase_report(parsed: dict, fmt: str, timestamp: str) -> dict:
    """导出测试用例报告"""
    try:
        if fmt == "md":
            md_lines = [f"# {parsed['title']}", ""]
            
            for case in parsed['cases']:
                md_lines.append(f"## {case['id']} {case['title']} [{case['priority']}]")
                md_lines.append("")
                md_lines.append("**测试步骤：**")
                for i, step in enumerate(case['steps'], 1):
                    md_lines.append(f"{i}. {step}")
                md_lines.append("")
                md_lines.append(f"**预期结果：**{case['expected']}")
                md_lines.append("")
            
            file_bytes = '\n'.join(md_lines).encode('utf-8')
            filename = f"testcase_report_{timestamp}.md"
        
        elif fmt == "pdf":
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
            from reportlab.lib import colors
            
            font_name = register_pdf_font()
            
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4,
                                  rightMargin=72, leftMargin=72,
                                  topMargin=72, bottomMargin=18)
            
            styles = getSampleStyleSheet()
            styles.add(ParagraphStyle(name='ChineseTitle', fontName=font_name, fontSize=18, spaceAfter=20))
            styles.add(ParagraphStyle(name='ChineseHeading', fontName=font_name, fontSize=14, spaceAfter=12, spaceBefore=12))
            styles.add(ParagraphStyle(name='ChineseSubHeading', fontName=font_name, fontSize=12, spaceAfter=8, spaceBefore=8, textColor=colors.HexColor('#5c4b37')))
            styles.add(ParagraphStyle(name='ChineseBody', fontName=font_name, fontSize=10, spaceAfter=6))
            
            story = []
            
            story.append(Paragraph(parsed['title'], styles['ChineseTitle']))
            story.append(Spacer(1, 0.2*inch))
            
            for case in parsed['cases']:
                story.append(Paragraph(f"{case['id']} {case['title']} [{case['priority']}]", styles['ChineseHeading']))
                
                story.append(Paragraph("测试步骤：", styles['ChineseSubHeading']))
                for i, step in enumerate(case['steps'], 1):
                    story.append(Paragraph(f"{i}. {step}", styles['ChineseBody']))
                
                story.append(Paragraph("预期结果：", styles['ChineseSubHeading']))
                story.append(Paragraph(case['expected'], styles['ChineseBody']))
                story.append(Spacer(1, 0.15*inch))
            
            doc.build(story)
            file_bytes = buffer.getvalue()
            filename = f"testcase_report_{timestamp}.pdf"
        
        elif fmt == "xlsx":
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
            
            wb = Workbook()
            ws = wb.active
            ws.title = "测试用例"
            
            # 标准测试用例表格格式
            headers = ['用例编号', '所属模块', '用例标题', '前置条件', '测试步骤', '预期结果', '优先级', '执行结果', '备注']
            ws.append(headers)
            
            # 表头样式 - 使用通用字体
            header_fill = PatternFill(start_color='C4A77D', end_color='C4A77D', fill_type='solid')
            header_font = Font(bold=True, color='FFFFFF', size=11, name='Arial')
            header_alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            thin_border = Border(
                left=Side(style='thin', color='d4c8b8'),
                right=Side(style='thin', color='d4c8b8'),
                top=Side(style='thin', color='d4c8b8'),
                bottom=Side(style='thin', color='d4c8b8')
            )
            
            for cell in ws[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = header_alignment
                cell.border = thin_border
            
            # 设置表头行高
            ws.row_dimensions[1].height = 30
            
            # 数据行样式 - 使用通用字体
            data_font = Font(size=10, name='Arial')
            data_alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
            priority_alignment = Alignment(horizontal='center', vertical='center')
            
            # 优先级颜色映射
            priority_colors = {
                'P0': 'FF6B6B',  # 红色
                'P1': 'FFB347',  # 橙色
                'P2': '87CEEB',  # 蓝色
            }
            
            for idx, case in enumerate(parsed['cases'], start=2):
                # 格式化测试步骤（带序号换行）
                steps_text = '\n'.join([f"{i}. {step}" for i, step in enumerate(case['steps'], 1)])
                
                # 提取模块信息（从标题中尝试提取）
                module = "登录功能"  # 默认模块，可根据实际情况调整
                
                row_data = [
                    case['id'],           # 用例编号
                    module,               # 所属模块
                    case['title'],        # 用例标题
                    "",                   # 前置条件（可扩展）
                    steps_text,           # 测试步骤
                    case['expected'],     # 预期结果
                    case['priority'],     # 优先级
                    "",                   # 执行结果（空白待填写）
                    ""                    # 备注
                ]
                ws.append(row_data)
                
                # 设置数据行样式
                for col_idx, cell in enumerate(ws[idx], start=1):
                    cell.font = data_font
                    cell.border = thin_border
                    
                    # 优先级列居中并着色
                    if col_idx == 7:  # 优先级列
                        cell.alignment = priority_alignment
                        priority = case['priority'].upper()
                        if priority in priority_colors:
                            cell.fill = PatternFill(start_color=priority_colors[priority], 
                                                   end_color=priority_colors[priority], 
                                                   fill_type='solid')
                            cell.font = Font(bold=True, color='FFFFFF', size=10, name='Arial')
                    # 用例编号居中
                    elif col_idx == 1:
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                    else:
                        cell.alignment = data_alignment
                
                # 设置行高（根据内容自动调整）
                ws.row_dimensions[idx].height = max(40, len(case['steps']) * 15 + 10)
            
            # 设置列宽
            ws.column_dimensions['A'].width = 12  # 用例编号
            ws.column_dimensions['B'].width = 15  # 所属模块
            ws.column_dimensions['C'].width = 30  # 用例标题
            ws.column_dimensions['D'].width = 25  # 前置条件
            ws.column_dimensions['E'].width = 45  # 测试步骤
            ws.column_dimensions['F'].width = 35  # 预期结果
            ws.column_dimensions['G'].width = 10  # 优先级
            ws.column_dimensions['H'].width = 12  # 执行结果
            ws.column_dimensions['I'].width = 20  # 备注
            
            # 冻结首行
            ws.freeze_panes = 'A2'
            
            buffer = BytesIO()
            wb.save(buffer)
            file_bytes = buffer.getvalue()
            filename = f"testcase_report_{timestamp}.xlsx"
        
        elif fmt == "docx":
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
            
            doc = Document()
            
            # 设置文档默认字体为通用字体
            def set_run_font(run, font_name='Arial', font_size=11, bold=False):
                run.font.name = font_name
                run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                run.font.size = Pt(font_size)
                run.bold = bold
            
            # 添加标题
            title = doc.add_heading(level=0)
            title_run = title.add_run(parsed['title'])
            set_run_font(title_run, font_size=18, bold=True)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            for case in parsed['cases']:
                h = doc.add_heading(level=2)
                set_run_font(h.add_run(f"{case['id']} {case['title']} [{case['priority']}]"), font_size=13, bold=True)
                
                h3 = doc.add_heading(level=3)
                set_run_font(h3.add_run("测试步骤："), font_size=11, bold=True)
                for i, step in enumerate(case['steps'], 1):
                    p = doc.add_paragraph(style='List Number')
                    set_run_font(p.add_run(f"{step}"), font_size=10)
                
                h3 = doc.add_heading(level=3)
                set_run_font(h3.add_run("预期结果："), font_size=11, bold=True)
                p = doc.add_paragraph()
                set_run_font(p.add_run(case['expected']), font_size=10)
                doc.add_paragraph()  # 空行
            
            buffer = BytesIO()
            doc.save(buffer)
            file_bytes = buffer.getvalue()
            filename = f"testcase_report_{timestamp}.docx"
        
        else:
            return {"success": False, "error": f"不支持的格式: {fmt}"}
        
        return {
            "success": True,
            "file": base64.b64encode(file_bytes).decode('utf-8'),
            "filename": filename
        }
    
    except Exception as e:
        import traceback
        return {"success": False, "error": f"导出失败: {str(e)}", "traceback": traceback.format_exc()}


@app.post("/export_single")
async def export_single(req: dict):
    """导出单条消息为报告格式"""
    key = get_api_key(req.get("session_token", ""))
    if not key:
        return {"success": False, "error": "未配置API Key"}
    
    # 优先使用 html_content（包含完整格式信息）
    # 如果没有 html_content，则使用 content（纯文本）
    html_content = req.get("html_content", "")
    text_content = req.get("content", "")
    content = html_content if html_content else text_content
    
    fmt = req.get("format", "md")
    bubble_type = req.get("bubble_type", "analysis")
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # 根据气泡类型选择解析和导出方式
    if bubble_type == "testcase":
        parsed = parse_testcase_content(content)
        return export_testcase_report(parsed, fmt, timestamp)
    else:
        # 默认需求分析
        parsed = parse_analysis_content(content)
        return export_analysis_report(parsed, fmt, timestamp)


@app.post("/test/run")
async def run_test():
    """运行Web自检（占位）"""
    return {"status": "developing", "message": "功能开发中"}

@app.get("/files/list")
async def list_files(path: str = "."):
    """获取文件列表（开发者模式）"""
    import os
    try:
        files = os.listdir(path)
        return {"files": files, "path": path}
    except Exception as e:
        return {"error": str(e)}

@app.get("/files/read")
async def read_file_api(path: str):
    """读取文件内容（开发者模式）"""
    try:
        with open(path, 'r') as f:
            content = f.read()
        return {"content": content, "path": path}
    except Exception as e:
        return {"error": str(e)}

@app.post("/files/save")
async def save_file_api(req: dict):
    """保存文件内容（开发者模式）"""
    try:
        path = req.get("path")
        content = req.get("content")
        with open(path, 'w') as f:
            f.write(content)
        return {"success": True, "path": path}
    except Exception as e:
        return {"error": str(e)}